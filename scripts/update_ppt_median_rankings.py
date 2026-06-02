#!/usr/bin/env python3
"""更新 PPT 中 Top 30 相关页面（池化聚合后）

聚合策略:
  E1 only: mean(4)
  E1+E2:   median(8)  — 跨评测池化
  E1+E3:   median(8)  — 选择性维度池化
  E1+E2+E3: median(12) — 全池化

更新页面：
  P24: Top 10 最终排名表 + 亮点摘要 + 聚合说明
  P41/A7: Top 30 三次评测排名对比（列标题更新）
  P44/A10: Top 30 论文基本信息表
  P45/A11: Top 30 评分详细表
"""

import csv
import json
import re
from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor


# ── Data Loading ──

def load_data():
    with open('results/unified_rankings.json') as f:
        uni = json.load(f)
    with open('/tmp/e1_mean_ranks.json') as f:
        e1_rank = {int(k): v for k, v in json.load(f).items()}
    with open('/tmp/e1e2_pooled_ranks.json') as f:
        e1e2_rank = {int(k): v for k, v in json.load(f).items()}

    meta = {}
    with open('results/merged-metadata.csv', 'r', encoding='utf-8-sig') as f:
        for row in csv.DictReader(f):
            meta[int(row['编号'])] = row

    return uni, e1_rank, e1e2_rank, meta


# ── Helper: set text preserving font ──

def set_text(shape, text, font_size=None, bold=None, color=None):
    """Set text on a shape, preserving existing font formatting."""
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    if tf.paragraphs and tf.paragraphs[0].runs:
        run = tf.paragraphs[0].runs[0]
        run.text = str(text)
        if font_size:
            run.font.size = font_size
        if bold is not None:
            run.font.bold = bold
        if color:
            run.font.color.rgb = color
    else:
        tf.paragraphs[0].text = str(text)


def get_textboxes_in_order(shapes):
    """Extract textboxes sorted by (top, left) position."""
    boxes = []
    for s in shapes:
        if s.has_text_frame and s.text_frame.text.strip():
            boxes.append(s)
    boxes.sort(key=lambda s: (s.top, s.left))
    return boxes


def find_group(shapes, name):
    """Find a group shape by name."""
    for s in shapes:
        if s.name == name:
            return s
    return None


# ── Color constants ──

RED = RGBColor(0xCC, 0x00, 0x00)
DARK_GRAY = RGBColor(0x1B, 0x1C, 0x21)  # default text color in PPT

# Source label colors (matching original PPT design)
SOURCE_COLORS = {
    'E1': RGBColor(0x5A, 0x64, 0x72),        # gray
    'E1+E2': RGBColor(0xC8, 0x16, 0x1E),     # red (was E2覆写)
    'E1+E3': RGBColor(0x2E, 0x7D, 0x32),     # green
    'E1+E2+E3': RGBColor(0x1E, 0x88, 0xE5),  # blue (new)
}

# Row background colors (light tints matching source colors)
BG_COLORS = {
    'E1': [RGBColor(0xFF, 0xFF, 0xFF), RGBColor(0xF8, 0xFA, 0xFC)],  # white / light gray alternating
    'E1+E2': RGBColor(0xFF, 0xEB, 0xEE),     # light red
    'E1+E3': RGBColor(0xE8, 0xF5, 0xE9),     # light green
    'E1+E2+E3': RGBColor(0xE3, 0xF2, 0xFD),  # light blue
}


# ── Source label helper ──

def source_label(item):
    """Generate display label for source."""
    src = item['source']
    if src == 'E1+E2+E3':
        return 'E1+E2+E3'
    elif src == 'E1+E2':
        return 'E1+E2'
    elif src == 'E1+E3':
        return 'E1+E3'
    else:
        return 'E1'


# ── P24: Top 10 Table ──

def update_p24(slide, top30, meta):
    """Update P24: Top 10 final ranking table + pipeline text + highlights."""

    # --- Update pipeline description (TextBox 9, outside group) ---
    for s in slide.shapes:
        if s.has_text_frame and 'E1 基线' in s.text_frame.text and 'E2' in s.text_frame.text:
            set_text(s, '1920 篇 E1 基线(均值) → 26 篇 E1+E2 池化(中位数) → 13 篇 E1+E3 池化(中位数) → 统一加权排名')
            break

    # --- Update table (Group 78) by TextBox name ---
    # Original PPT uses fixed TextBox numbering:
    #   Header: TextBox 13-17
    #   Row 1: TextBox 19-23 (rank, title, journal, score, source)
    #   Row 2: TextBox 25-29, ..., Row 10: TextBox 73-77
    # Pattern: row i (1-indexed) → TextBox (19 + (i-1)*6) to (23 + (i-1)*6)
    group78 = find_group(slide.shapes, 'Group 78')
    if not group78:
        print('  WARNING: Group 78 not found on P24')
        return

    # Build name→shape map
    name_map = {}
    for s in group78.shapes:
        if s.has_text_frame:
            name_map[s.name] = s

    for i in range(10):
        item = top30[i]
        pid = item['pid']
        m = meta.get(pid, {})
        title = m.get('题目', f'Paper {pid}')[:25]
        journal = m.get('期刊', '?')
        score = f'{item["weighted_score"]:.2f}'
        src = source_label(item)

        rank_label = f'★{i+1}' if item['source'] != 'E1' else str(i+1)

        base_num = 19 + i * 6
        rank_box = name_map.get(f'TextBox {base_num}')
        title_box = name_map.get(f'TextBox {base_num + 1}')
        journal_box = name_map.get(f'TextBox {base_num + 2}')
        score_box = name_map.get(f'TextBox {base_num + 3}')
        source_box = name_map.get(f'TextBox {base_num + 4}')

        if rank_box:
            set_text(rank_box, rank_label)
        if title_box:
            set_text(title_box, title)
        if journal_box:
            set_text(journal_box, journal)
        if score_box:
            set_text(score_box, score)
        if source_box:
            set_text(source_box, src)

    # --- Update highlights (Group 83) ---
    group83 = find_group(slide.shapes, 'Group 83')
    if group83:
        e2_count = sum(1 for item in top30 if item['e2_override'])
        e3_count = sum(1 for item in top30 if item['e3_merged'])
        e1_count = len(top30) - e2_count - e3_count

        for s in group83.shapes:
            if not s.has_text_frame:
                continue
            text = s.text_frame.text.strip()
            if '排名变化亮点' in text:
                set_text(s, '排名变化亮点（中位数池化聚合后）')
            elif '新冠军' in text:
                set_text(s,
                    f'新冠军: pid-1260 公司决议瑕疵 (88.53, E1); '
                    f'E1+E2 池化 {e2_count} 篇; E1+E3 池化 {e3_count} 篇; '
                    f'E1 基线 {e1_count} 篇'
                )
            elif 'E2 覆写' in text or 'E1+E2 池化' in text:
                e2_pids = [str(item['pid']) for item in top30 if item['e2_override']]
                e3_pids = [f'#{item["rank"]}' for item in top30 if item['e3_merged']]
                set_text(s,
                    f'E1+E2 池化 {len(e2_pids)} 篇 (pid {", ".join(e2_pids[:4])}...); '
                    f'E1+E3 池化 {len(e3_pids)} 篇 ({", ".join(e3_pids[:4])}...); '
                    f'为什么中位数: 4模型均值易受单模型极端值影响, 8+评分取中位数可自动过滤'
                )


# ── P41/A7: Top 30 Ranking Comparison ──

def update_p41(slide, top30, e1_rank, e1e2_rank, meta):
    """Update P41: 30-row E1/E1+E2/最终 ranking comparison."""

    # --- Update column headers ---
    for s in slide.shapes:
        if not s.has_text_frame:
            continue
        text = s.text_frame.text.strip()
        if text == 'E2':
            set_text(s, 'E1+E2')
        elif text == 'E3最终':
            set_text(s, '最终')
        elif text == 'Δ(E3-E1)':
            set_text(s, 'Δ')
        elif text == '附录 A7：Top 30 三次评测排名对比':
            set_text(s, '附录 A7：Top 30 排名对比（中位数池化聚合）')
        elif 'E1/E2/E3 展示' in text:
            set_text(s, 'E1=全量评审均值排名(4模型); E1+E2=E1与E2评分池化后中位数排名(8评分); 最终=E1+E2+E3池化。中位数池化可自动过滤单模型极端值。')

    # --- Update data rows ---
    boxes = get_textboxes_in_order(slide.shapes)

    header_idx = None
    for i, b in enumerate(boxes):
        if b.text_frame.text.strip() == '#':
            header_idx = i
            break

    if header_idx is None:
        print('  WARNING: Header "#" not found on P41')
        return

    data_start = header_idx + 6

    # Update footnote
    for b in boxes:
        text = b.text_frame.text.strip()
        if '说明：' in text:
            set_text(b, '说明：E1=全量评审均值排名(4模型); E1+E2=池化中位数排名(8评分); 最终=E1+E2池化+E3选择性维度池化。绿色行=池化导致排名显著变化。')

    for row in range(30):
        idx = data_start + row * 6
        if idx + 5 >= len(boxes):
            break

        item = top30[row]
        pid = item['pid']
        m = meta.get(pid, {})
        title = m.get('题目', f'Paper {pid}')[:22]

        e1r = e1_rank.get(pid, '?')

        if item['e2_override']:
            e2r = str(e1e2_rank.get(pid, '?'))
        else:
            e2r = '—'

        final_r = str(item['rank'])

        if isinstance(e1r, int):
            delta = e1r - item['rank']
            delta_str = f'{delta:+d}' if delta != 0 else '0'
        else:
            delta_str = '?'

        set_text(boxes[idx], str(row + 1))
        set_text(boxes[idx + 1], title)
        set_text(boxes[idx + 2], str(e1r))
        set_text(boxes[idx + 3], e2r)
        set_text(boxes[idx + 4], final_r)
        set_text(boxes[idx + 5], delta_str)


# ── P44/A10: Top 30 Basic Info ──

def update_p44(slide, top30, meta):
    """Update P44: 30-row paper basic info table."""
    boxes = get_textboxes_in_order(slide.shapes)

    header_idx = None
    for i, b in enumerate(boxes):
        if b.text_frame.text.strip() == '排名':
            header_idx = i
            break

    if header_idx is None:
        print('  WARNING: Header "排名" not found on P44')
        return

    data_start = header_idx + 7

    # Update title to mention pooled aggregation
    for b in boxes:
        text = b.text_frame.text.strip()
        if '附录 A10' in text:
            set_text(b, '附录 A10：Top 30 论文基本信息表（中位数池化聚合）')

    # Count sources for footnote
    source_counts = {}
    for item in top30:
        src = source_label(item)
        source_counts[src] = source_counts.get(src, 0) + 1

    # Update footnote
    for b in boxes:
        text = b.text_frame.text.strip()
        if '数据源' in text and 'E1基线' in text:
            e1_count = source_counts.get('E1', 0)
            e1e2_count = source_counts.get('E1+E2', 0)
            e1e3_count = source_counts.get('E1+E3', 0)
            e1e2e3_count = source_counts.get('E1+E2+E3', 0)
            set_text(b, f'数据源：E1基线({e1_count}篇) + E1+E2池化({e1e2_count}篇) + E1+E3池化({e1e3_count}篇)；E1+E2+E3池化({e1e2e3_count}篇)。机构核验：Top30 无上海交大论文。')

    for row in range(30):
        idx = data_start + row * 7
        if idx + 6 >= len(boxes):
            break

        item = top30[row]
        pid = item['pid']
        m = meta.get(pid, {})

        title = m.get('题目', f'Paper {pid}')[:25]
        journal = m.get('期刊', '?')
        year = m.get('年份', '?')
        author = m.get('作者', '?')
        if author:
            author = re.sub(r'\[[\d,\s]+\]', '', author)
            if ';' in author:
                author = author.split(';')[0].strip()
        institution = m.get('作者机构', '?')
        if institution and ';' in institution:
            institution = institution.split(';')[0]
        if institution:
            institution = re.sub(r'\[[\d,\s]+\]', '', institution).strip()

        set_text(boxes[idx], str(row + 1))
        set_text(boxes[idx + 1], title)
        set_text(boxes[idx + 2], journal)
        set_text(boxes[idx + 3], str(year))
        set_text(boxes[idx + 4], author if author else '?')
        set_text(boxes[idx + 5], institution if institution else '?')

        # Source label with color coding
        src = source_label(item)
        src_color = SOURCE_COLORS.get(src, DARK_GRAY)
        set_text(boxes[idx + 6], src, color=src_color)

    # ── Update row background rectangles ──
    # Collect source label positions and their corresponding source types
    src_positions = []
    for s in slide.shapes:
        if s.has_text_frame:
            text = s.text_frame.text.strip()
            if text in SOURCE_COLORS:
                src_positions.append((s.top, text))

    # Collect row background rectangles (skip header/decoration rects with small top)
    bg_rects = []
    for s in slide.shapes:
        if 'Rectangle' in s.name and s.top > 1200000:
            try:
                _ = s.fill.fore_color.rgb  # check if fillable
                bg_rects.append(s)
            except:
                pass

    bg_rects.sort(key=lambda s: s.top)

    # Match each rectangle to the nearest source label and set color
    e1_toggle = 0
    for rect in bg_rects:
        # Find nearest source label by top position
        best_src = None
        best_dist = float('inf')
        for src_top, src_label in src_positions:
            dist = abs(rect.top - src_top)
            if dist < best_dist:
                best_dist = dist
                best_src = src_label

        if best_src and best_dist < 80000:  # tolerance for matching
            bg = BG_COLORS.get(best_src)
            if isinstance(bg, list):
                # E1: alternate white/gray
                rect.fill.solid()
                rect.fill.fore_color.rgb = bg[e1_toggle % 2]
                e1_toggle += 1
            elif bg:
                rect.fill.solid()
                rect.fill.fore_color.rgb = bg


# ── P45/A11: Top 30 Scoring Details ──

def update_p45(slide, top30, meta):
    """Update P45: 30-row scoring details table."""
    boxes = get_textboxes_in_order(slide.shapes)

    header_idx = None
    for i, b in enumerate(boxes):
        if b.text_frame.text.strip() == '#':
            header_idx = i
            break

    if header_idx is None:
        print('  WARNING: Header "#" not found on P45')
        return

    data_start = header_idx + 9

    # Count sources for footnote
    source_counts = {}
    for item in top30:
        src = source_label(item)
        source_counts[src] = source_counts.get(src, 0) + 1

    # Update title and footnote
    for b in boxes:
        text = b.text_frame.text.strip()
        if '附录 A11' in text:
            set_text(b, '附录 A11：Top 30 评分详细表（中位数池化聚合）')
        elif '红色 = 标准差' in text:
            set_text(b, '红色 = 标准差 > 5，需专家复核。聚合方式: E1=均值(4模型); E1+E2/E3=中位数(8评分池化)。中位数可自动过滤单模型极端值，避免排名被异常评分绑架。')
        elif '数据源' in text and 'E1 基线' in text:
            e1_count = source_counts.get('E1', 0)
            e1e2_count = source_counts.get('E1+E2', 0)
            e1e3_count = source_counts.get('E1+E3', 0)
            e1e2e3_count = source_counts.get('E1+E2+E3', 0)
            set_text(b, f'数据源：E1 基线(均值) + E1+E2 池化({e1e2_count}篇) + E1+E3 池化({e1e3_count}篇)；E1+E2+E3 池化({e1e2e3_count}篇)；红色 = std > 5，需专家复核。')

    DIM_KEYS = [
        'problem_originality',
        'literature_insight',
        'analytical_framework',
        'logical_coherence',
        'conclusion_consensus',
        'forward_extension',
    ]

    for row in range(30):
        idx = data_start + row * 9
        if idx + 8 >= len(boxes):
            break

        item = top30[row]
        pid = item['pid']
        m = meta.get(pid, {})
        title = m.get('题目', f'Paper {pid}')[:20]

        ws = item['weighted_score']
        wstd = item['weighted_std']
        total_str = f'{ws:.2f}±{wstd:.2f}'

        set_text(boxes[idx], str(row + 1))
        set_text(boxes[idx + 1], title)
        set_text(boxes[idx + 2], total_str)

        for d, dim_key in enumerate(DIM_KEYS):
            avg = item['dim_avgs'].get(dim_key, 0)
            std = item['dim_stds'].get(dim_key, 0)
            # Round first, then check — so display matches the red decision
            avg_r = round(avg, 1)
            std_r = round(std, 1)
            dim_str = f'{avg_r:.1f}±{std_r:.1f}'
            box = boxes[idx + 3 + d]
            if std_r > 5:
                set_text(box, dim_str, bold=False, color=RED)
            else:
                set_text(box, dim_str, bold=False, color=DARK_GRAY)


# ── Main ──

def main():
    uni, e1_rank, e1e2_rank, meta = load_data()
    top30 = uni['top30']

    pptx_path = 'projects/sjtu-socialeval-report_ppt169_20260601/exports/sjtu-socialeval-report_20260602_222842.pptx'
    prs = Presentation(pptx_path)

    slides = list(prs.slides)

    print('Updating P24 (Top 10 table)...')
    update_p24(slides[23], top30, meta)

    print('Updating P41 (A7: ranking comparison)...')
    update_p41(slides[40], top30, e1_rank, e1e2_rank, meta)

    print('Updating P44 (A10: basic info)...')
    update_p44(slides[43], top30, meta)

    print('Updating P45 (A11: scoring details)...')
    update_p45(slides[44], top30, meta)

    print('Updating P32 (核心成果总结)...')
    update_p32(slides[31], top30)

    out_path = 'projects/sjtu-socialeval-report_ppt169_20260601/exports/sjtu-socialeval-report_20260603_pooled.pptx'
    prs.save(out_path)
    print(f'\nSaved to: {out_path}')

def update_p32(slide, top30):
    """Update P32: 核心成果总结 with pooled aggregation info."""

    # Build name→shape map
    name_map = {}
    for s in slide.shapes:
        if s.has_text_frame:
            name_map[s.name] = s

    # Count sources
    source_counts = {}
    for item in top30:
        src = source_label(item)
        source_counts[src] = source_counts.get(src, 0) + 1

    # Update E1 description (TextBox 12)
    if 'TextBox 12' in name_map:
        set_text(name_map['TextBox 12'], '均值(4模型)形成基线排名')

    # Update E2 description (TextBox 16, 17)
    if 'TextBox 16' in name_map:
        set_text(name_map['TextBox 16'], '26 篇与 E1 池化')
    if 'TextBox 17' in name_map:
        set_text(name_map['TextBox 17'], '8 评分取中位数，自动过滤极端值')

    # Update E3 description (TextBox 21, 22)
    if 'TextBox 21' in name_map:
        e3_count = source_counts.get('E1+E3', 0) + source_counts.get('E1+E2+E3', 0)
        set_text(name_map['TextBox 21'], f'{e3_count} 篇不稳定维度')
    if 'TextBox 22' in name_map:
        set_text(name_map['TextBox 22'], '与现有评分池化取中位数')

    # Update Top30 composition (TextBox 26)
    if 'TextBox 26' in name_map:
        e1_count = source_counts.get('E1', 0)
        e1e2_count = source_counts.get('E1+E2', 0)
        e1e3_count = source_counts.get('E1+E3', 0)
        set_text(name_map['TextBox 26'], f'E1:{e1_count} + E1+E2:{e1e2_count} + E1+E3:{e1e3_count}')

    # Update key methodology (TextBox 28)
    if 'TextBox 28' in name_map:
        set_text(name_map['TextBox 28'],
            '关键口径：E1=均值(基线)，E1+E2/E3=中位数(池化)。中位数自动过滤单模型极端值，避免排名被异常评分绑架。')


if __name__ == '__main__':
    main()
