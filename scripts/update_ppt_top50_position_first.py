#!/usr/bin/env python3
"""Update the 0603 SocialEval deck to the position-first Top50 rule."""

from __future__ import annotations

import json
import shutil
from pathlib import Path
from typing import Iterable

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
SOURCE = (
    ROOT
    / "projects"
    / "sjtu-socialeval-report_ppt169_20260601"
    / "exports"
    / "中国法学自主知识体系创新AI辅助评价汇报0603.pptx"
)
OUTPUT = (
    ROOT
    / "projects"
    / "sjtu-socialeval-report_ppt169_20260601"
    / "exports"
    / "中国法学自主知识体系创新AI辅助评价汇报0603_top50_position_first.pptx"
)
TOP50 = ROOT / "results" / "top101" / "top50-position-first-proportional.json"

FONT = "PingFang SC"
FONT_ALT = "Microsoft YaHei"
BG = RGBColor(248, 246, 239)
INK = RGBColor(35, 36, 38)
MUTED = RGBColor(96, 99, 105)
RED = RGBColor(132, 28, 42)
RED_DARK = RGBColor(94, 20, 31)
GOLD = RGBColor(194, 145, 59)
GREEN = RGBColor(50, 113, 91)
LINE = RGBColor(218, 211, 197)
WHITE = RGBColor(255, 255, 255)


def load_top50() -> dict:
    with TOP50.open("r", encoding="utf-8") as fh:
        return json.load(fh)


def clear_slide(slide) -> None:
    for shape in list(slide.shapes):
        slide.shapes._spTree.remove(shape._element)


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    size: float = 16,
    color: RGBColor = INK,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_rect(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    fill: RGBColor = WHITE,
    line: RGBColor | None = LINE,
):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.6)
    return shape


def setup_slide(slide, number: int, title: str, subtitle: str = "") -> None:
    clear_slide(slide)
    add_rect(slide, 0, 0, 13.333, 7.5, BG, None)
    add_rect(slide, 0, 0, 0.16, 7.5, RED, None)
    add_text(slide, f"{number}", 0.3, 0.18, 0.35, 0.25, 8, RED, True)
    add_text(slide, title, 0.65, 0.15, 10.8, 0.45, 21, RED_DARK, True)
    if subtitle:
        add_text(slide, subtitle, 0.68, 0.62, 11.8, 0.32, 9.5, MUTED)
    add_rect(slide, 0.65, 0.98, 11.95, 0.015, GOLD, None)


def bullet_box(slide, title: str, bullets: Iterable[str], x: float, y: float, w: float, h: float):
    add_rect(slide, x, y, w, h, WHITE, LINE)
    add_text(slide, title, x + 0.18, y + 0.12, w - 0.36, 0.25, 12.5, RED_DARK, True)
    text = "\n".join(f"• {item}" for item in bullets)
    add_text(slide, text, x + 0.2, y + 0.45, w - 0.4, h - 0.55, 9.2, INK)


def metric(slide, label: str, value: str, note: str, x: float, y: float, w: float):
    add_rect(slide, x, y, w, 1.0, WHITE, LINE)
    add_text(slide, value, x + 0.15, y + 0.15, w - 0.3, 0.32, 20, RED_DARK, True)
    add_text(slide, label, x + 0.15, y + 0.50, w - 0.3, 0.20, 8.8, INK, True)
    add_text(slide, note, x + 0.15, y + 0.73, w - 0.3, 0.20, 7.4, MUTED)


def simple_table(slide, rows: list[list[str]], x: float, y: float, w: float, h: float, font_size: float):
    table_shape = slide.shapes.add_table(len(rows), len(rows[0]), Inches(x), Inches(y), Inches(w), Inches(h))
    table = table_shape.table
    for r_idx, row in enumerate(rows):
        for c_idx, value in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = str(value)
            cell.margin_left = Inches(0.02)
            cell.margin_right = Inches(0.02)
            cell.margin_top = Inches(0.01)
            cell.margin_bottom = Inches(0.01)
            fill = RED_DARK if r_idx == 0 else (WHITE if r_idx % 2 else RGBColor(252, 250, 245))
            cell.fill.solid()
            cell.fill.fore_color.rgb = fill
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.LEFT
                for run in paragraph.runs:
                    run.font.name = FONT_ALT
                    run.font.size = Pt(font_size)
                    run.font.bold = r_idx == 0
                    run.font.color.rgb = WHITE if r_idx == 0 else INK
    return table_shape


def truncate(text: str, length: int) -> str:
    return text if len(text) <= length else text[: length - 1] + "…"


def update_slide_4(slide, top50: dict) -> None:
    setup_slide(slide, 4, "评价对象与规模", "从 1920 篇全量评审到 Top101 稳定性复测，再到归属优先的 Top50 配额榜")
    metric(slide, "可评测论文总量", "1920", "近十年法学三大刊", 0.75, 1.25, 2.25)
    metric(slide, "Top101 候选池", "101", "Top60 + 年份/学科覆盖", 3.25, 1.25, 2.25)
    metric(slide, "五轴归属评估", "101/101", "两模型 + 条件 R2", 5.75, 1.25, 2.25)
    metric(slide, "正式 Top50", "50", "归属优先 + 学科比例", 8.25, 1.25, 2.25)
    metric(slide, "五轴 10 分", "49", "主资格池入选", 10.75, 1.25, 1.6)
    bullet_box(
        slide,
        "当前评测成果",
        [
            "E1：1920 篇全量 R1+R2 评审完成，形成六维质量基线。",
            "E2：Top101 进入稳定性复测，覆盖每年至少 5 篇、各学科至少 5 篇。",
            "五轴：101 篇完成中国法学自主知识体系位置归属度评估。",
            "Top50：按 1920 全库学科比例配额，从五轴 10 分池中择优；仅 1 篇 9 分用于学科补足。",
        ],
        0.75,
        2.65,
        5.8,
        3.6,
    )
    bullet_box(
        slide,
        "Top50 新口径",
        [
            "先判归属资格：10 分为主池，9 分只用于配额不足补足。",
            "再按学科配额：沿用全库 1920 篇学科比例。",
            "最后按六维质量排序：学科内按加权创新分择优。",
            "8 分及以下不进入正式 Top50，保留为观察或专家复核样本。",
        ],
        6.8,
        2.65,
        5.55,
        3.6,
    )


def update_slide_6(slide) -> None:
    setup_slide(slide, 6, "四阶段评审流程", "五轴位置归属度采用 0-10 分结构化评价，不进入六维基础分")
    stages = [
        ("阶段一", "预检 · 项目口径判断", "进入评分 / 边界复核 / 明显不适格"),
        ("阶段二", "六维评分 · 学术质量评价", "0-100 分；构成六维加权基础分"),
        ("阶段三", "五轴位置归属度评价", "对象、材料、范畴、解释目标、体系映射；0-10 分"),
        ("阶段四", "评价层复核 · 可靠性判断", "标准差、模型分歧、证据缺口与专家终审"),
    ]
    for i, (tag, title, note) in enumerate(stages):
        x = 0.75 + i * 3.05
        add_rect(slide, x, 1.55, 2.55, 3.95, WHITE, LINE)
        add_text(slide, tag, x + 0.15, 1.75, 2.2, 0.25, 10, GOLD, True, PP_ALIGN.CENTER)
        add_text(slide, title, x + 0.18, 2.25, 2.2, 0.62, 15, RED_DARK, True, PP_ALIGN.CENTER)
        add_text(slide, note, x + 0.22, 3.35, 2.12, 1.2, 10, INK, False, PP_ALIGN.CENTER)
        if i < 3:
            add_text(slide, "→", x + 2.65, 3.25, 0.35, 0.35, 22, RED, True, PP_ALIGN.CENTER)
    add_text(slide, "核心原则：五轴归属只决定候选资格与知识体系画像，不作为第七个质量维度，不进入六维加权分。", 0.9, 6.3, 11.2, 0.35, 12, RED_DARK, True, PP_ALIGN.CENTER)


def update_slide_8(slide) -> None:
    setup_slide(slide, 8, "概念操作化：把主观判断变成可复核证据", "六维评价判断学术质量；五轴归属判断知识单元能否进入中国法学自主知识体系")
    cards = [
        ("枢纽性", "影响范围、争议焦点、理论基础、制度连锁等可检查标准。"),
        ("可争辩法学问题", "必须有问题句、法律规范关联、理论分歧，并区别于政策倡导。"),
        ("原创性理论", "看新命题、系统论证框架、对既有理论的可定位突破。"),
        ("五轴位置归属度", "对象、材料、范畴、解释目标、体系映射五轴各 0/1/2 分。"),
    ]
    for i, (title, body) in enumerate(cards):
        x = 0.75 + (i % 2) * 6.0
        y = 1.35 + (i // 2) * 2.35
        add_rect(slide, x, y, 5.55, 1.75, WHITE, LINE)
        add_text(slide, f"0{i+1}", x + 0.16, y + 0.14, 0.45, 0.28, 11, GOLD, True)
        add_text(slide, title, x + 0.72, y + 0.13, 4.55, 0.32, 15, RED_DARK, True)
        add_text(slide, body, x + 0.72, y + 0.62, 4.55, 0.78, 10.2, INK)
    add_text(slide, "意义：AI 不直接裁定“好不好”，而是把问题、证据、分歧和归属位置拆成可检查字段，供专家复核。", 0.85, 6.25, 11.4, 0.35, 12, RED_DARK, True, PP_ALIGN.CENTER)


def update_slide_13(slide) -> None:
    setup_slide(slide, 13, "五轴位置归属度评价", "回答“论文产出的知识单元能否纳入中国法学自主知识体系”，不替代六维质量评价")
    axes = [
        ("对象归属度", "核心问题是否属于中国法秩序、制度、实践、争论或传统。"),
        ("材料归属度", "核心材料是否来自中国规范、案例、制度、史料、实务或中文争论。"),
        ("范畴自主度", "核心分析范畴是否经过中国法语境重置，而非直接套用外部框架。"),
        ("解释目标归属度", "最终解释目标是否指向中国法学知识生产。"),
        ("体系映射度", "知识产出能否映射到既有节点、细化节点、交叉节点或候选新增节点。"),
    ]
    for i, (title, body) in enumerate(axes):
        x = 0.75 + (i % 3) * 4.0
        y = 1.28 + (i // 3) * 2.05
        add_rect(slide, x, y, 3.55, 1.55, WHITE, LINE)
        add_text(slide, str(i + 1), x + 0.12, y + 0.10, 0.34, 0.30, 13, GOLD, True, PP_ALIGN.CENTER)
        add_text(slide, title, x + 0.52, y + 0.10, 2.75, 0.28, 12.2, RED_DARK, True)
        add_text(slide, body, x + 0.18, y + 0.55, 3.15, 0.75, 8.8, INK)
    rows = [
        ["分值", "含义"],
        ["0", "无明确证据，或只是背景性提及"],
        ["1", "有局部证据，但不是论文核心结构"],
        ["2", "构成核心问题、材料、范畴、解释目标或明确体系位置"],
    ]
    simple_table(slide, rows, 0.9, 5.65, 11.2, 1.05, 9.2)


def update_slide_14(slide) -> None:
    setup_slide(slide, 14, "五轴防饱和与复核触发", "避免把“研究中国问题”机械等同于强自主知识体系归属")
    bullet_box(
        slide,
        "防饱和规则",
        [
            "对象、材料、解释目标可共享证据，但不能因为“讨论中国法”自动三轴全给 2。",
            "常规中国法条解释不能自动推出范畴自主度=2 或体系映射度=2。",
            "范畴自主度=2 必须说明核心范畴如何由中国制度、实践、传统或法学争论生成。",
            "体系映射度=2 必须给出明确既有节点、细化节点、交叉节点或候选新增节点。",
        ],
        0.75,
        1.3,
        5.75,
        4.55,
    )
    bullet_box(
        slide,
        "复核触发",
        [
            "高归属分但证据薄弱：进入专家复核。",
            "候选新增节点：优先专家确认，关系到知识树扩展。",
            "模型在范畴自主度或体系映射度上分歧明显：保留分歧解释。",
            "高六维分但归属低：不否定质量，转为高质量一般法学论文或观察样本。",
        ],
        6.8,
        1.3,
        5.55,
        4.55,
    )
    add_text(slide, "结论：五轴是“归属资格与知识谱系画像”，不是论文质量的第七个评分维度。", 0.9, 6.25, 11.2, 0.4, 12, RED_DARK, True, PP_ALIGN.CENTER)


def update_slide_15(slide) -> None:
    setup_slide(slide, 15, "知识库匹配与学科手册", "体系映射度优先复用中国法学自主知识体系树状知识库，但匹配数量不等于质量分")
    metric(slide, "知识库学科", "23", "含主干与新兴方向", 0.8, 1.25, 2.5)
    metric(slide, "标识性概念", "333", "用于节点候选", 3.55, 1.25, 2.5)
    metric(slide, "原创性理论", "255", "用于理论谱系", 6.3, 1.25, 2.5)
    metric(slide, "框架结构", "92", "用于结构映射", 9.05, 1.25, 2.5)
    bullet_box(
        slide,
        "知识库提供什么",
        [
            "既有节点匹配：学科、概念、理论、框架结构。",
            "交叉节点提示：跨部门法、跨范畴的知识位置。",
            "候选新增节点：有证据但知识树尚未覆盖的表达。",
        ],
        0.8,
        2.75,
        5.45,
        2.65,
    )
    bullet_box(
        slide,
        "不能替代什么",
        [
            "不能替代五轴评分：知识树可能不完整。",
            "不能替代六维质量评价：匹配多不等于论文更好。",
            "不能替代专家确认：候选新增节点需人工把关。",
        ],
        6.7,
        2.75,
        5.45,
        2.65,
    )
    add_text(slide, "Top50 新名单已更新；知识库匹配统计需按同一 50 篇名单重跑，历史口径数字不再混用。", 0.9, 6.15, 11.2, 0.45, 11.5, RED_DARK, True, PP_ALIGN.CENTER)


def update_slide_16(slide) -> None:
    setup_slide(slide, 16, "大规模评测概览", "E1 全量评审、E2 Top101 稳定性复测、E3 选择性补测共同支撑 Top50 归属优先配额榜")
    cards = [
        ("E1 全量评审", "1920 篇", "形成 R1+R2 六维质量基线，约 92,160 次 API 调用。"),
        ("E2 Top101 复测", "101 篇", "Top60 + 年份至少 5 篇 + 学科至少 5 篇。"),
        ("E3 选择性补测", "45 篇", "63 个不稳定维度进入选择性补测。"),
        ("五轴归属评估", "101 篇", "两模型评估，必要时条件 R2。"),
        ("Top50 输出", "50 篇", "10 分主池 + 学科比例配额；1 篇 9 分补足。"),
    ]
    for i, (title, value, note) in enumerate(cards):
        x = 0.75 + (i % 3) * 4.0
        y = 1.28 + (i // 3) * 2.15
        add_rect(slide, x, y, 3.55, 1.65, WHITE, LINE)
        add_text(slide, value, x + 0.18, y + 0.16, 3.15, 0.38, 19, RED_DARK, True)
        add_text(slide, title, x + 0.18, y + 0.60, 3.15, 0.22, 9, INK, True)
        add_text(slide, note, x + 0.18, y + 0.92, 3.15, 0.45, 8.2, MUTED)
    add_text(slide, "主评测 API 调用仍以 E1/E2/E3 为主；五轴归属度用于 Top50 资格和知识体系画像，不改变六维质量分。", 0.9, 6.22, 11.2, 0.4, 11.5, RED_DARK, True, PP_ALIGN.CENTER)


def update_slide_22(slide) -> None:
    setup_slide(slide, 22, "E2 入口：Top101 稳定性复测池", "从纯高分候选扩展为高分 + 年份覆盖 + 学科覆盖，避免候选池过度集中")
    bullet_box(
        slide,
        "入口规则",
        [
            "加权分 Top60 形成高质量基线。",
            "年份覆盖：每年至少 5 篇，避免早期年份被挤出。",
            "学科覆盖：各学科至少 5 篇，避免民商法/刑法/法理过度集中。",
            "最终 E2 候选池：101 篇，作为 Top50 的上游候选真源。",
        ],
        0.75,
        1.35,
        5.6,
        4.65,
    )
    rows = [
        ["年份", "篇数"],
        ["2015", "6"],
        ["2016", "8"],
        ["2017", "6"],
        ["2018", "7"],
        ["2019", "11"],
        ["2020", "9"],
        ["2021", "11"],
        ["2022", "11"],
        ["2023", "12"],
        ["2024", "10"],
        ["2025", "10"],
    ]
    simple_table(slide, rows, 6.75, 1.35, 2.35, 4.55, 8.2)
    rows2 = [["学科覆盖", "篇数"], ["刑法学", "20"], ["民商法学", "19"], ["法学理论", "10"], ["诉讼法学", "9"], ["宪行", "8"], ["其他 7 学科", "各 5"]]
    simple_table(slide, rows2, 9.35, 1.35, 2.8, 2.75, 8.2)
    add_text(slide, "Top101 负责覆盖性；Top50 在此基础上先做五轴归属资格筛选，再按全库学科比例配额输出。", 6.75, 4.55, 5.4, 0.8, 11, RED_DARK, True)


def update_slide_23(slide) -> None:
    setup_slide(slide, 23, "E2 结果：Top101 稳定性复测", "稳定性分析看同一模型、同一论文、同一维度的 E1 R2 → E2 R2 差值")
    metric(slide, "E2 候选池", "101", "全部进入 E1/E2 合并排名", 0.8, 1.25, 2.6)
    metric(slide, "来源分布", "56 + 45", "56 篇 E1+E2；45 篇 E1+E2+E3", 3.7, 1.25, 2.6)
    metric(slide, "五轴完成率", "101/101", "用于 Top50 资格判断", 6.6, 1.25, 2.6)
    metric(slide, "Top50 入选", "50", "配额榜分数范围 89.225-82.600", 9.5, 1.25, 2.6)
    bullet_box(
        slide,
        "如何使用 E2",
        [
            "E2 不只是重排 Top 分数，而是检验同一模型重复评测波动。",
            "E1+E2 使用多源分数池化，中位数降低单模型极端值影响。",
            "E3 只补测不稳定维度，不再全量重复调用。",
        ],
        0.8,
        2.8,
        5.55,
        2.9,
    )
    bullet_box(
        slide,
        "如何进入 Top50",
        [
            "先通过五轴归属资格：10 分主池，9 分仅补足。",
            "再按 1920 全库学科比例确定配额。",
            "最后在学科内按六维加权分选取。",
        ],
        6.7,
        2.8,
        5.55,
        2.9,
    )


def update_slide_24(slide) -> None:
    setup_slide(slide, 24, "E3：选择性补测", "只补测不稳定维度，降低成本，同时保留专家复核线索")
    metric(slide, "选择性补测论文", "45", "E3 selective pool", 0.8, 1.25, 2.55)
    metric(slide, "不稳定维度", "63", "进入 E3 的维度单元", 3.65, 1.25, 2.55)
    metric(slide, "调用规模", "~504", "完整 R1+R2 估算调用", 6.5, 1.25, 2.55)
    metric(slide, "用途", "复核线索", "不改变五轴归属规则", 9.35, 1.25, 2.55)
    bullet_box(
        slide,
        "选择性补测原则",
        [
            "只针对高分候选中的不稳定维度，而不是重新全量评测。",
            "E3 结果进入 E1/E2/E3 池化，提高高分候选排序稳健性。",
            "仍无法收敛的维度不强行平均，保留为专家终审线索。",
        ],
        0.8,
        2.85,
        5.55,
        2.75,
    )
    bullet_box(
        slide,
        "与 Top50 的关系",
        [
            "六维加权分用于学科内排序。",
            "五轴归属分用于 Top50 资格和补足判断。",
            "标准差高不自动剔除，但在报告中标记需复核。",
        ],
        6.7,
        2.85,
        5.55,
        2.75,
    )


def update_slide_25(slide, papers: list[dict]) -> None:
    setup_slide(slide, 25, "Top50 最终排名：归属优先 + 学科比例配额", "本页展示新规则前 10 名；完整名单见附录")
    rows = [["#", "PID", "学科", "论文", "六维分", "五轴"]]
    for paper in papers[:10]:
        rows.append([
            str(paper["rank"]),
            str(paper["pid"]),
            paper["category"],
            truncate(paper["title"], 28),
            f'{paper["score"]:.3f}',
            str(paper["position_total_score"]),
        ])
    simple_table(slide, rows, 0.65, 1.25, 12.05, 4.65, 8.2)
    add_text(slide, "唯一口径：results/top101/top50-position-first-proportional.json / .csv", 0.85, 6.15, 11.5, 0.25, 10, MUTED)
    add_text(slide, "分数范围 89.225-82.600；49 篇五轴 10 分，1 篇五轴 9 分用于宪法学与行政法学配额补足。", 0.85, 6.45, 11.5, 0.3, 11, RED_DARK, True)


def update_slide_26(slide, metadata: dict) -> None:
    setup_slide(slide, 26, "Top50 五轴归属与知识库匹配口径", "当前版先固化 Top50 入选名单；知识库匹配统计必须按同一 50 篇名单重跑")
    metric(slide, "Top50 总数", "50", "归属优先配额榜", 0.8, 1.25, 2.6)
    metric(slide, "五轴 10 分", "49", "主资格池", 3.7, 1.25, 2.6)
    metric(slide, "五轴 9 分", "1", "学科配额补足", 6.6, 1.25, 2.6)
    metric(slide, "≤8 分", "0", "不入正式 Top50", 9.5, 1.25, 2.6)
    bullet_box(
        slide,
        "本版可确认",
        [
            "Top50 排名、学科配额、六维分、标准差和五轴分均来自新 canonical 数据。",
            "新名单较旧配额版替换 6 篇，清除了五轴低于 9 分的旧入选论文。",
            "五轴用于资格与画像，不进入六维加权分。",
        ],
        0.8,
        2.85,
        5.55,
        2.85,
    )
    bullet_box(
        slide,
        "不混用旧数字",
        [
            "历史位置归属均值、条目数和高置信度比例不再展示。",
            "既有知识库匹配结果对应旧名单，不直接迁移到新名单。",
            "新知识库匹配应以 `top50-position-first-proportional.json` 的 50 个 PID 为输入重跑。",
        ],
        6.7,
        2.85,
        5.55,
        2.85,
    )


def update_slide_30(slide) -> None:
    setup_slide(slide, 30, "评测流程全景图", "候选识别 → Top101 稳定性复测 → E3 选择性补测 → Top50 归属优先配额榜")
    steps = [
        ("E1 全量评审", "1920 篇生成 R1+R2，形成六维质量基线。"),
        ("E2 Top101 复测", "Top60 + 年份覆盖 + 学科覆盖，观察重复评测稳定性。"),
        ("E3 选择性补测", "只补测不稳定维度，避免全量重复调用。"),
        ("五轴归属评估", "对象、材料、范畴、解释目标、体系映射五轴 0-10。"),
        ("Top50 配额榜", "10 分主池 + 全库学科比例配额 + 学科内六维分排序。"),
    ]
    for i, (title, note) in enumerate(steps):
        y = 1.25 + i * 0.95
        add_rect(slide, 1.2, y, 10.7, 0.62, WHITE, LINE)
        add_text(slide, f"{i+1}", 1.38, y + 0.12, 0.35, 0.2, 10, GOLD, True, PP_ALIGN.CENTER)
        add_text(slide, title, 1.9, y + 0.10, 2.4, 0.22, 11, RED_DARK, True)
        add_text(slide, note, 4.25, y + 0.10, 7.25, 0.24, 9.5, INK)
    add_text(slide, "A6/A7 展示的是候选形成与稳定性路径；Top50 是归属资格和学科配额后的最终展示清单。", 1.0, 6.25, 11.3, 0.36, 11, RED_DARK, True, PP_ALIGN.CENTER)


def update_slide_31(slide) -> None:
    setup_slide(slide, 31, "核心成果总结", "方法论流程保持不变，候选池与最终展示名单升级为 Top101 + 五轴归属 + Top50 配额榜")
    bullet_box(
        slide,
        "方法论创新",
        [
            "六维评价：研究创新性、现状洞察度、理论建构力、逻辑连贯性、学术共识度、前瞻延展性。",
            "五轴归属：对象、材料、范畴、解释目标、体系映射，用于自主知识体系位置判断。",
            "交叉评审：R1 独立评分 → R2 交叉审视，暴露真实学术分歧。",
        ],
        0.75,
        1.25,
        5.55,
        4.65,
    )
    bullet_box(
        slide,
        "工程成果",
        [
            "E1：1920 篇全量 R1+R2 评审完成。",
            "E2：101 篇完成稳定性复测，覆盖年份与学科。",
            "五轴：101 篇完成位置归属度评估。",
            "Top50：按全库学科比例输出，49 篇 10 分、1 篇 9 分补足。",
        ],
        6.75,
        1.25,
        5.55,
        4.65,
    )
    add_text(slide, "核心指标：Top50 分数范围 89.225-82.600；正式名单不含五轴 8 分及以下论文。", 0.9, 6.25, 11.2, 0.35, 12, RED_DARK, True, PP_ALIGN.CENTER)


def update_slide_32(slide) -> None:
    setup_slide(slide, 32, "局限与下一步", "下一步重点是五轴归属证据、知识树节点和专家复核")
    bullet_box(
        slide,
        "当前局限",
        [
            "部分高分论文仍存在研究创新性或理论建构力分歧，需专家终审。",
            "五轴中的范畴自主度、体系映射度最容易产生可解释分歧。",
            "新 Top50 知识库匹配需要按同一 50 篇名单重跑，不能复用历史统计。",
        ],
        0.75,
        1.35,
        5.75,
        4.6,
    )
    bullet_box(
        slide,
        "下一步工作",
        [
            "专家复核：保留标准差高、证据薄弱或候选新增节点的论文。",
            "知识树扩展：对候选新增节点进行专家确认。",
            "报告输出：主榜、配额说明、五轴证据和六维分歧并列展示。",
            "学科拓展：从法学扩展到人文社科其他学科。",
        ],
        6.8,
        1.35,
        5.55,
        4.6,
    )


def update_slide_34(slide) -> None:
    setup_slide(slide, 34, "附录 A1：评分公式详解", "六维加权分用于质量排序；五轴归属度用于候选资格与知识体系画像")
    add_text(slide, "base_score = 创新性×30% + 洞察度×20% + 建构力×15% + 连贯性×20% + 共识度×10% + 延展性×5%", 0.9, 1.35, 11.3, 0.45, 14, RED_DARK, True, PP_ALIGN.CENTER)
    bullet_box(
        slide,
        "六维计分原则",
        [
            "六维评分各自独立，避免互相补偿。",
            "权重反映“命题与定位”的优先性。",
            "总分只作为候选排序线索，不替代专家终审。",
        ],
        0.85,
        2.25,
        5.45,
        3.2,
    )
    bullet_box(
        slide,
        "五轴归属说明",
        [
            "五轴位置归属度 = 0-10 分，独立于六维质量分。",
            "Top50 先看五轴资格，再看学科配额，最后按六维分排序。",
            "高质量 + 强归属 + 明确知识节点 = 自主知识体系代表性候选论文。",
        ],
        6.8,
        2.25,
        5.45,
        3.2,
    )


def update_slide_35(slide) -> None:
    setup_slide(slide, 35, "附录 A2：典型论文评审报告示例", "单篇报告由六维质量、五轴归属、模型分歧和复核建议四部分构成")
    bullet_box(
        slide,
        "输出内容",
        [
            "六维均分与标准差：展示每个维度的质量判断与分歧。",
            "五轴归属证据：对象、材料、范畴、解释目标、体系映射。",
            "知识树节点：既有节点、交叉节点或候选新增节点。",
            "模型意见：保留 accepted / rejected points，解释为什么分歧。",
        ],
        0.75,
        1.3,
        5.75,
        4.65,
    )
    bullet_box(
        slide,
        "专家复核入口",
        [
            "任一维度 std > 5：标记分歧，必要时专家复核。",
            "候选新增节点：专家确认是否进入知识树。",
            "五轴高分但证据不足：复核归属判断，而不是直接否定论文质量。",
            "专家评分与 AI 评分并列存储，不互相覆盖。",
        ],
        6.8,
        1.3,
        5.55,
        4.65,
    )


def update_slide_39(slide, metadata: dict) -> None:
    setup_slide(slide, 39, "附录 A6：Top50 新旧名单差异", "新规则清除五轴低于 9 分的旧入选论文，并补入同学科高归属候选")
    removed = metadata["removed_from_previous_proportional"]
    added = metadata["added_vs_previous_proportional"]
    rows = [["移除旧 PID", "原因"], *[[str(pid), "五轴归属低于 9 分，不符合正式 Top50 资格"] for pid in removed]]
    simple_table(slide, rows, 0.75, 1.35, 5.65, 3.0, 8.0)
    rows2 = [["补入新 PID", "原因"], *[[str(pid), "五轴归属 10 分，同学科按六维分顺延入选"] for pid in added]]
    simple_table(slide, rows2, 6.75, 1.35, 5.65, 3.0, 8.0)
    add_text(slide, "差异结果：移除 6 篇，补入 6 篇；新名单仍满足原 1920 全库学科比例配额。", 0.9, 5.05, 11.2, 0.36, 12, RED_DARK, True, PP_ALIGN.CENTER)
    add_text(slide, "规则解释：若某学科 10 分论文不足，才允许 9 分补足；8 分及以下不入正式 Top50。", 0.9, 5.55, 11.2, 0.36, 11.5, INK, False, PP_ALIGN.CENTER)


def update_slide_41(slide, metadata: dict) -> None:
    setup_slide(slide, 41, "附录 A8：Top50 五轴分布与入选来源", "五轴归属度用于资格筛选，六维加权分用于学科内排序")
    metric(slide, "Top50", "50", "正式名单总数", 0.8, 1.25, 2.45)
    metric(slide, "10 分主池", "49", "98.0%", 3.55, 1.25, 2.45)
    metric(slide, "9 分补足", "1", "2.0%", 6.3, 1.25, 2.45)
    metric(slide, "≤8 分", "0", "不入选", 9.05, 1.25, 2.45)
    rows = [["学科", "配额"], *[[k, str(v)] for k, v in metadata["discipline_quotas"].items()]]
    simple_table(slide, rows, 0.8, 2.75, 4.0, 3.8, 7.4)
    bullet_box(
        slide,
        "入选逻辑",
        [
            "先限定五轴 10 分主池。",
            "按全库学科比例分配 50 个名额。",
            "学科内按六维创新加权分排序。",
            "入选来源含：五轴归属9分学科配额补足 1 篇。",
        ],
        5.2,
        2.75,
        6.65,
        3.8,
    )


def info_rows(papers: list[dict]) -> list[list[str]]:
    rows = [["#", "学科", "题名", "期刊/年", "作者", "分数", "五轴"]]
    for p in papers:
        rows.append([
            str(p["rank"]),
            p["category"].replace("宪法学与行政法学", "宪行"),
            truncate(p["title"], 24),
            f'{p["journal"]}/{p["year"]}',
            truncate(p["author"], 8),
            f'{p["score"]:.2f}',
            str(p["position_total_score"]),
        ])
    return rows


def score_rows(papers: list[dict]) -> list[list[str]]:
    rows = [["#", "论文", "总分±std", "创新", "洞察", "建构", "连贯", "共识", "延展", "五轴"]]
    for p in papers:
        d = p["dimensions"]
        def val(key: str) -> str:
            item = d[key]
            return f'{float(item["pooled_avg"]):.1f}±{float(item["pooled_std"]):.1f}'
        rows.append([
            str(p["rank"]),
            truncate(p["title"], 18),
            f'{p["score"]:.2f}±{float(p["std"]):.1f}',
            val("problem_originality"),
            val("literature_insight"),
            val("analytical_framework"),
            val("logical_coherence"),
            val("conclusion_consensus"),
            val("forward_extension"),
            str(p["position_total_score"]),
        ])
    return rows


def update_info_slide(slide, number: int, title: str, papers: list[dict]) -> None:
    setup_slide(slide, number, title, "完整 Top50 基本信息表，按新规则最终排名排序")
    simple_table(slide, info_rows(papers), 0.35, 1.15, 12.65, 5.85, 5.7)


def update_score_slide(slide, number: int, title: str, papers: list[dict]) -> None:
    setup_slide(slide, number, title, "格式：均分±std；标准差高的条目保留专家复核线索")
    simple_table(slide, score_rows(papers), 0.25, 1.15, 12.85, 5.85, 4.7)


def add_blank_slide(prs):
    # The 0603 template only exposes four layouts; use the title-only layout
    # and immediately clear it in setup_slide().
    layout = prs.slide_layouts[2]
    return prs.slides.add_slide(layout)


def main() -> None:
    top50 = load_top50()
    papers = top50["papers"]
    metadata = top50["metadata"]

    shutil.copyfile(SOURCE, OUTPUT)
    prs = Presentation(str(OUTPUT))

    updates = {
        4: lambda s: update_slide_4(s, top50),
        6: update_slide_6,
        8: update_slide_8,
        13: update_slide_13,
        14: update_slide_14,
        15: update_slide_15,
        16: update_slide_16,
        22: update_slide_22,
        23: update_slide_23,
        24: update_slide_24,
        25: lambda s: update_slide_25(s, papers),
        26: lambda s: update_slide_26(s, metadata),
        30: update_slide_30,
        31: update_slide_31,
        32: update_slide_32,
        34: update_slide_34,
        35: update_slide_35,
        39: lambda s: update_slide_39(s, metadata),
        41: lambda s: update_slide_41(s, metadata),
        42: lambda s: update_info_slide(s, 42, "附录 A9：Top50 论文基本信息表（1-25）", papers[:25]),
        43: lambda s: update_info_slide(s, 43, "附录 A10：Top50 论文基本信息表（26-50）", papers[25:]),
    }

    for slide_number, updater in updates.items():
        updater(prs.slides[slide_number - 1])

    slide44 = add_blank_slide(prs)
    update_score_slide(slide44, 44, "附录 A11：Top50 评分详细表（1-25）", papers[:25])
    slide45 = add_blank_slide(prs)
    update_score_slide(slide45, 45, "附录 A12：Top50 评分详细表（26-50）", papers[25:])

    prs.save(str(OUTPUT))
    print(OUTPUT)


if __name__ == "__main__":
    main()
