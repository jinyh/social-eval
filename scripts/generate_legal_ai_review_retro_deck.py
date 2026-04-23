from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUT_DIR = Path("slide-deck/legal-ai-review-retro")
PPTX_PATH = OUT_DIR / "legal-ai-review-retro.pptx"

BG = RGBColor(243, 246, 250)
PANEL = RGBColor(255, 255, 255)
TEXT = RGBColor(15, 23, 42)
MUTED = RGBColor(71, 85, 105)
BLUE = RGBColor(47, 93, 140)
CYAN = RGBColor(88, 166, 199)
GOLD = RGBColor(203, 170, 92)
LINE = RGBColor(203, 213, 225)

FONT_HEAD = "PingFang SC"
FONT_BODY = "PingFang SC"


def add_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG

    # Top header strip
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.35))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BLUE
    shape.line.fill.background()

    # Bottom accent strip
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(7.15), Inches(13.333), Inches(0.35))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(229, 236, 242)
    shape.line.fill.background()


def add_panel(slide, left, top, width, height, *, fill=PANEL, line=LINE):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1.2)
    return shape


def add_text(slide, left, top, width, height, text, *, size=20, bold=False, color=TEXT, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    font = run.font
    font.name = FONT_HEAD if bold else FONT_BODY
    font.size = Pt(size)
    font.bold = bold
    font.color.rgb = color
    return box


def add_bullets(slide, left, top, width, items, *, size=18, color=TEXT, bullet_color=BLUE, line_gap=0.34):
    y = top
    row_height = max(0.5, line_gap + 0.14)
    for item in items:
        add_text(slide, left, y, Inches(0.22), Inches(row_height - 0.08), "•", size=size, bold=True, color=bullet_color)
        add_text(slide, left + Inches(0.26), y - Inches(0.01), width - Inches(0.26), Inches(row_height), item, size=size, color=color)
        y += Inches(row_height)


def add_title(slide, title, subtitle=None):
    add_text(slide, Inches(0.65), Inches(0.55), Inches(12.0), Inches(0.7), title, size=28, bold=True, color=TEXT)
    if subtitle:
        add_text(slide, Inches(0.68), Inches(1.2), Inches(11.8), Inches(0.45), subtitle, size=14, color=MUTED)


def add_chip(slide, left, top, width, text, *, fill=RGBColor(233, 242, 248), color=BLUE):
    chip = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, Inches(0.42))
    chip.fill.solid()
    chip.fill.fore_color.rgb = fill
    chip.line.fill.background()
    add_text(slide, left + Inches(0.12), top + Inches(0.06), width - Inches(0.2), Inches(0.22), text, size=11, bold=True, color=color, align=PP_ALIGN.CENTER)


def add_footer_note(slide, text):
    add_text(slide, Inches(0.72), Inches(7.02), Inches(12.0), Inches(0.18), text, size=9, color=MUTED)


def connector(slide, x1, y1, x2, y2, color=BLUE, width=2):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(width)
    return line


def slide_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    hero = add_panel(slide, Inches(0.7), Inches(0.95), Inches(11.9), Inches(4.8), fill=RGBColor(248, 250, 252))
    hero.line.color.rgb = RGBColor(193, 208, 222)
    add_text(slide, Inches(1.05), Inches(1.45), Inches(10.8), Inches(1.3), "法学论文 AI 辅助评审逻辑与流程复盘", size=28, bold=True)
    add_text(slide, Inches(1.07), Inches(2.55), Inches(7.8), Inches(0.5), "面向专家意见征询与机制共建", size=16, color=BLUE)
    add_text(slide, Inches(1.08), Inches(3.18), Inches(9.6), Inches(0.8), "这不是某个版本的发布说明，而是一套正在收敛中的法学评审机制复盘。", size=18)
    add_text(slide, Inches(1.08), Inches(3.78), Inches(9.6), Inches(0.8), "目标是请专家帮助校准规则、指出边界、支持后续试点推进。", size=18)

    # abstract six-node motif
    centers = [(9.75, 2.15), (10.6, 2.75), (10.6, 3.75), (9.75, 4.35), (8.9, 3.75), (8.9, 2.75)]
    for x, y in centers:
        node = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(y), Inches(0.32), Inches(0.32))
        node.fill.solid()
        node.fill.fore_color.rgb = CYAN
        node.line.fill.background()
    for i in range(len(centers)):
        x1, y1 = centers[i]
        x2, y2 = centers[(i + 1) % len(centers)]
        connector(slide, Inches(x1 + 0.16), Inches(y1 + 0.16), Inches(x2 + 0.16), Inches(y2 + 0.16), CYAN, 1.5)

    badge = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(9.05), Inches(1.2), Inches(2.25), Inches(0.46))
    badge.fill.solid()
    badge.fill.fore_color.rgb = GOLD
    badge.line.fill.background()
    add_text(slide, Inches(9.18), Inches(1.29), Inches(2.0), Inches(0.2), "专家征询版", size=12, bold=True, color=RGBColor(99, 66, 12), align=PP_ALIGN.CENTER)

    add_chip(slide, Inches(1.05), Inches(5.15), Inches(1.78), "评审标准显化")
    add_chip(slide, Inches(2.95), Inches(5.15), Inches(1.95), "判断链条可复核")
    add_chip(slide, Inches(5.02), Inches(5.15), Inches(2.18), "邀请专家参与校准")
    add_footer_note(slide, "用途：专家研讨 / 机制评审 / 共建沟通")


def slide_problem(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "传统评审的难点不只在效率，更在标准难显化", "专家判断高度依赖隐性经验，难以直接迁移给 AI")
    add_panel(slide, Inches(0.72), Inches(1.7), Inches(6.15), Inches(4.8))
    add_text(slide, Inches(1.0), Inches(1.98), Inches(3.0), Inches(0.28), "三个真正的结构性难点", size=18, bold=True, color=BLUE)
    card_specs = [
        (2.42, 0.86, "标准隐含于共同体经验", "很多判断依赖共同体默认知识，不总能直接写成规则。", 15, 12),
        (3.48, 0.86, "解释过程难复用", "同样的判断理由难以沉淀为可复核、可迁移的稳定过程。", 15, 12),
        (4.54, 1.02, "AI 只能模仿语言，不能自动模仿判断", "如果不先外化标准，模型容易生成像评审的话，却未必真的按评审逻辑判断。", 13, 11),
    ]
    for y, height, title, body, title_size, body_size in card_specs:
        card = add_panel(slide, Inches(0.98), Inches(y), Inches(5.55), Inches(height), fill=RGBColor(250, 252, 255))
        card.line.color.rgb = RGBColor(214, 223, 232)
        add_chip(slide, Inches(1.15), Inches(y + 0.16), Inches(1.55), "关键难点", fill=RGBColor(233, 242, 248))
        title_height = 0.22 if height <= 0.9 else 0.34
        body_top = y + 0.42 if height <= 0.9 else y + 0.5
        body_height = 0.26 if height <= 0.9 else 0.34
        add_text(slide, Inches(2.88), Inches(y + 0.12), Inches(2.95), Inches(title_height), title, size=title_size, bold=True, color=TEXT)
        add_text(slide, Inches(2.88), Inches(body_top), Inches(3.18), Inches(body_height), body, size=body_size, color=MUTED)
    tri = [(8.35, 2.05), (10.95, 3.15), (8.35, 4.55)]
    labels = ["标准显化", "评审效率", "解释可复用"]
    colors = [BLUE, GOLD, CYAN]
    for (x, y), label, color in zip(tri, labels, colors):
        node = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(y), Inches(1.7), Inches(0.72))
        node.fill.solid()
        node.fill.fore_color.rgb = RGBColor(246, 249, 252)
        node.line.color.rgb = color
        node.line.width = Pt(2)
        add_text(slide, Inches(x + 0.18), Inches(y + 0.2), Inches(1.35), Inches(0.22), label, size=14, bold=True, color=color, align=PP_ALIGN.CENTER)
    connector(slide, Inches(9.2), Inches(2.77), Inches(10.95), Inches(3.5), GOLD, 1.5)
    connector(slide, Inches(9.2), Inches(4.91), Inches(10.95), Inches(3.88), CYAN, 1.5)
    connector(slide, Inches(9.2), Inches(2.77), Inches(9.2), Inches(4.55), BLUE, 1.5)
    add_footer_note(slide, "判断困难的核心不是工作量，而是专家标准长期以隐性知识形式存在。")


def slide_workflow(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "当前主链路把论文评审拆成可复核的七个环节", "从论文输入到报告输出，每一步都对应明确的判断任务")
    steps = [
        "文档预处理", "前置检查", "六维评分", "可靠性判断", "报告生成", "专家复核接口", "版本留存",
    ]
    x = 0.72
    for i, step in enumerate(steps):
        color = BLUE if i < 5 else (GOLD if i == 5 else CYAN)
        panel = add_panel(slide, Inches(x), Inches(2.55), Inches(1.6), Inches(1.38), fill=RGBColor(250, 252, 255))
        panel.line.color.rgb = color
        add_text(slide, Inches(x + 0.12), Inches(2.72), Inches(1.35), Inches(0.5), f"{i+1:02d}", size=12, bold=True, color=color)
        add_text(slide, Inches(x + 0.12), Inches(3.03), Inches(1.35), Inches(0.55), step, size=13, bold=True)
        if i < len(steps) - 1:
            connector(slide, Inches(x + 1.6), Inches(3.25), Inches(x + 1.8), Inches(3.25), color, 2)
        x += 1.8
    add_chip(slide, Inches(1.05), Inches(5.0), Inches(1.55), "规则驱动")
    add_chip(slide, Inches(2.75), Inches(5.0), Inches(1.55), "模型输出", fill=RGBColor(232, 247, 251), color=CYAN)
    add_chip(slide, Inches(4.45), Inches(5.0), Inches(1.75), "人工介入", fill=RGBColor(249, 243, 229), color=RGBColor(125, 94, 35))
    add_bullets(slide, Inches(7.75), Inches(4.85), Inches(4.8), [
        "主链路依据当前配置、需求文档与代码实现整理。",
        "目标不是自动定结论，而是拆解判断节点。",
    ], size=15, line_gap=0.4)
    add_footer_note(slide, "主链路依据：当前 YAML 配置、需求文档 v0.4、评审主流程代码。")


def slide_intent_gap(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "AI 常常听懂了词面，却没有准确对齐评审意图", "同样的评审要求，不同模型并不会自动进入同一个法学语境")
    add_panel(slide, Inches(0.9), Inches(2.0), Inches(2.8), Inches(2.7), fill=RGBColor(247, 250, 253))
    add_text(slide, Inches(1.15), Inches(2.25), Inches(2.2), Inches(0.35), "专家评审意图", size=18, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    add_bullets(slide, Inches(1.14), Inches(2.85), Inches(2.15), [
        "创新性如何定义",
        "洞察度何为充分",
        "可接受性边界在哪",
    ], size=15, line_gap=0.42)
    add_panel(slide, Inches(4.55), Inches(1.85), Inches(3.95), Inches(3.0), fill=RGBColor(251, 253, 255))
    add_text(slide, Inches(5.18), Inches(2.08), Inches(2.7), Inches(0.25), "YAML 规则层", size=18, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    add_chip(slide, Inches(5.0), Inches(2.58), Inches(1.08), "定义")
    add_chip(slide, Inches(6.17), Inches(2.58), Inches(1.08), "顺序")
    add_chip(slide, Inches(7.34), Inches(2.58), Inches(1.08), "上限")
    add_chip(slide, Inches(5.0), Inches(3.15), Inches(1.08), "证据")
    add_chip(slide, Inches(6.17), Inches(3.15), Inches(1.08), "输出")
    add_chip(slide, Inches(7.34), Inches(3.15), Inches(1.08), "复核")
    for y, label in [(1.95, "模型 A"), (3.15, "模型 B"), (4.35, "模型 C")]:
        card = add_panel(slide, Inches(9.35), Inches(y), Inches(2.8), Inches(0.95), fill=RGBColor(250, 252, 255))
        card.line.color.rgb = CYAN
        add_text(slide, Inches(9.58), Inches(y + 0.12), Inches(2.25), Inches(0.25), label, size=15, bold=True, color=CYAN)
        add_text(slide, Inches(9.58), Inches(y + 0.42), Inches(2.0), Inches(0.18), "输出解释并不天然一致", size=12, color=MUTED)
    connector(slide, Inches(3.7), Inches(3.35), Inches(4.55), Inches(3.35), BLUE, 2)
    connector(slide, Inches(8.5), Inches(2.35), Inches(9.35), Inches(2.35), CYAN, 2)
    connector(slide, Inches(8.5), Inches(3.6), Inches(9.35), Inches(3.6), CYAN, 2)
    connector(slide, Inches(8.5), Inches(4.8), Inches(9.35), Inches(4.8), CYAN, 2)
    add_footer_note(slide, "规则必须从“隐含理解”转成“显式控制”，否则输出会随模型和时间漂移。")


def slide_yaml(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "YAML 不是“参数表”，而是评审机制真正的规则层", "它把专家经验从隐含判断变成可加载、可执行、可修改的规则集合")
    add_panel(slide, Inches(0.8), Inches(1.95), Inches(3.3), Inches(3.9), fill=RGBColor(250, 252, 255)).line.color.rgb = BLUE
    add_text(slide, Inches(1.02), Inches(2.2), Inches(2.6), Inches(0.3), "它不是只写六维和权重", size=19, bold=True, color=BLUE)
    add_bullets(slide, Inches(1.04), Inches(2.78), Inches(2.7), [
        "还会写 precheck。",
        "还会写评分锚定与封顶规则。",
        "还会写输出格式与复核标记。",
    ], size=15, line_gap=0.5)

    center = add_panel(slide, Inches(4.45), Inches(1.7), Inches(4.4), Inches(4.45), fill=RGBColor(248, 251, 254))
    center.line.color.rgb = CYAN
    add_text(slide, Inches(4.9), Inches(2.05), Inches(3.5), Inches(0.35), "YAML 规则层", size=24, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
    chips = ["precheck", "decision_order", "scoring_criteria", "ceiling_rules", "output_contract", "review_flags"]
    positions = [(5.0, 2.85), (6.62, 2.85), (5.0, 3.55), (6.62, 3.55), (5.0, 4.25), (6.62, 4.25)]
    for label, (x, y) in zip(chips, positions):
        add_chip(slide, Inches(x), Inches(y), Inches(1.35), label, fill=RGBColor(245, 249, 253))
    add_text(slide, Inches(5.0), Inches(5.0), Inches(3.4), Inches(0.7), "从专家视角看，它就是把评审经验写成一套可审阅、可修改、可追踪的规则接口。", size=15, color=MUTED, align=PP_ALIGN.CENTER)

    add_panel(slide, Inches(9.1), Inches(1.95), Inches(3.45), Inches(3.9), fill=RGBColor(250, 252, 255)).line.color.rgb = GOLD
    add_text(slide, Inches(9.33), Inches(2.2), Inches(2.7), Inches(0.3), "接下来三页分别解释", size=19, bold=True, color=RGBColor(125, 94, 35))
    add_bullets(slide, Inches(9.36), Inches(2.8), Inches(2.7), [
        "它先定义“评什么”。",
        "再定义“怎么判”。",
        "最后定义“怎么输出与复核”。",
    ], size=15, bullet_color=GOLD, line_gap=0.56)
    add_footer_note(slide, "这页的作用：先把 YAML 从“技术细节”提升为“机制中心”。")


def slide_yaml_scope(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "YAML 先定义“评什么”，并规定判断顺序", "先把评价对象和程序写清，模型才不会直接跳到结论")
    left = add_panel(slide, Inches(0.82), Inches(1.9), Inches(4.15), Inches(4.95), fill=RGBColor(250, 252, 255))
    left.line.color.rgb = BLUE
    add_text(slide, Inches(1.08), Inches(2.16), Inches(2.6), Inches(0.3), "这一层写入的对象", size=18, bold=True, color=BLUE)
    rows = [
        ("precheck", "决定是否进入六维评分"),
        ("dimensions", "定义六维结构"),
        ("weight", "定义总分计算关系"),
        ("decision_order", "规定先看什么后看什么"),
    ]
    y = 2.82
    for key, desc in rows:
        add_chip(slide, Inches(1.1), Inches(y), Inches(1.28), key, fill=RGBColor(245, 249, 253))
        add_text(slide, Inches(2.58), Inches(y + 0.08), Inches(1.95), Inches(0.25), desc, size=14)
        y += 0.72

    right = add_panel(slide, Inches(5.35), Inches(1.9), Inches(6.85), Inches(4.95), fill=RGBColor(250, 252, 255))
    right.line.color.rgb = CYAN
    add_text(slide, Inches(5.65), Inches(2.16), Inches(3.0), Inches(0.3), "运行时的对应作用", size=18, bold=True, color=CYAN)
    add_bullets(slide, Inches(5.68), Inches(2.72), Inches(5.85), [
        "先跑 precheck，判断论文是 pass、conditional_pass，还是 reject。",
        "再按 dimensions 逐维评分，而不是让模型一次性凭整体印象给总分。",
        "最后按 weight 计算综合得分，并保留每个维度独立判断。",
        "decision_order 强制模型按固定程序判断，避免跳步和偷看结论。",
    ], size=15, bullet_color=CYAN, line_gap=0.5)
    add_footer_note(slide, "这一层控制的是：评价对象和评价程序。")


def slide_yaml_rules(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "YAML 再定义“怎么判”，也就是锚定区间与封顶边界", "真正影响分数稳定性的，不只是 prompt，而是 scoring_criteria 和 ceiling_rules")
    left = add_panel(slide, Inches(0.85), Inches(1.95), Inches(5.15), Inches(4.8), fill=RGBColor(250, 252, 255))
    left.line.color.rgb = BLUE
    add_text(slide, Inches(1.15), Inches(2.2), Inches(2.9), Inches(0.3), "评分锚定是怎么工作的", size=18, bold=True, color=BLUE)
    bands = [
        ("excellent", "高位区间，要求关键条件基本成立", BLUE),
        ("good", "条件成立但强度有限", CYAN),
        ("marginal", "勉强成立，问题明显", GOLD),
        ("unacceptable", "基础条件不成立", RGBColor(155, 86, 74)),
    ]
    y = 2.9
    for name, desc, color in bands:
        chip = add_chip(slide, Inches(1.18), Inches(y), Inches(1.2), name, fill=RGBColor(245, 249, 253), color=color)
        add_text(slide, Inches(2.58), Inches(y + 0.08), Inches(2.95), Inches(0.25), desc, size=14)
        y += 0.68

    right = add_panel(slide, Inches(6.28), Inches(1.95), Inches(5.3), Inches(4.8), fill=RGBColor(250, 252, 255))
    right.line.color.rgb = GOLD
    add_text(slide, Inches(6.58), Inches(2.2), Inches(3.2), Inches(0.3), "封顶和边界为什么重要", size=18, bold=True, color=RGBColor(125, 94, 35))
    add_bullets(slide, Inches(6.62), Inches(2.82), Inches(4.45), [
        "ceiling_rules 决定：哪些基础条件不成立时，分数必须封顶。",
        "它防止模型在“问题并未成立”或“研究地图明显残缺”时仍然给出高分。",
        "boundary_notes 会提醒：本维度评什么、不评什么，减少维度串位。",
        "这一步控制的是：分档依据和封顶边界。",
    ], size=15, bullet_color=GOLD, line_gap=0.52)
    add_footer_note(slide, "这页解释的是：为什么真正左右分数稳定性的，往往是规则而不是语言润色。")


def slide_yaml_output_hooks(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "YAML 最后定义“怎么输出”以及“何时提示复核”", "结果不仅要能生成，还要能被比较、核验和人工接手")
    left = add_panel(slide, Inches(0.82), Inches(1.95), Inches(3.75), Inches(4.85), fill=RGBColor(250, 252, 255))
    left.line.color.rgb = CYAN
    add_text(slide, Inches(1.08), Inches(2.2), Inches(2.6), Inches(0.3), "结果生成前", size=18, bold=True, color=CYAN)
    add_bullets(slide, Inches(1.12), Inches(2.84), Inches(2.9), [
        "prompt_template 决定模型按什么路径进入该维度。",
        "它会把本维度的定义、判断步骤、字数限制、输出格式要求一起带进去。",
    ], size=15, bullet_color=CYAN, line_gap=0.58)

    mid = add_panel(slide, Inches(4.8), Inches(1.95), Inches(3.55), Inches(4.85), fill=RGBColor(250, 252, 255))
    mid.line.color.rgb = BLUE
    add_text(slide, Inches(5.06), Inches(2.2), Inches(2.6), Inches(0.3), "结果生成时", size=18, bold=True, color=BLUE)
    add_bullets(slide, Inches(5.1), Inches(2.84), Inches(2.7), [
        "output_contract 约束必须返回哪些字段。",
        "缺字段、乱字段、分数不服从规则，都应该视为失败输出。",
    ], size=15, bullet_color=BLUE, line_gap=0.58)

    right = add_panel(slide, Inches(8.6), Inches(1.95), Inches(4.0), Inches(4.85), fill=RGBColor(250, 252, 255))
    right.line.color.rgb = GOLD
    add_text(slide, Inches(8.88), Inches(2.2), Inches(3.0), Inches(0.3), "结果生成后", size=18, bold=True, color=RGBColor(125, 94, 35))
    add_bullets(slide, Inches(8.92), Inches(2.84), Inches(3.1), [
        "review_flags 把证据不足、边界模糊、需要外部核验等情况显式标出。",
        "这样专家不是接手一团文本，而是接手一个已标注风险点的结构化结果。",
    ], size=15, bullet_color=GOLD, line_gap=0.58)
    add_footer_note(slide, "这一层控制的是：结果是否可读、可比、可复核，以及何时应提醒人工接手。")


def slide_dimensions(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "六维框架把法学论文的判断链条外化出来", "从起点到落脚，六个维度共同构成结构化评价路径")
    center_x, center_y = 4.0, 3.55
    items = [
        ("问题创新性\n30%", 4.0, 1.95, BLUE),
        ("现状洞察度\n15%", 5.6, 2.7, CYAN),
        ("分析框架建构力\n15%", 5.6, 4.1, CYAN),
        ("逻辑严密性\n25%", 4.0, 4.85, BLUE),
        ("结论可接受性\n10%", 2.4, 4.1, GOLD),
        ("前瞻延展性\n5%", 2.4, 2.7, GOLD),
    ]
    for label, x, y, color in items:
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.HEXAGON, Inches(x), Inches(y), Inches(1.45), Inches(0.92))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(250, 252, 255)
        shape.line.color.rgb = color
        shape.line.width = Pt(2)
        add_text(slide, Inches(x + 0.08), Inches(y + 0.18), Inches(1.28), Inches(0.5), label, size=12, bold=True, color=color, align=PP_ALIGN.CENTER)
    for _, x, y, color in items:
        connector(slide, Inches(center_x + 0.72), Inches(center_y + 0.45), Inches(x + 0.72), Inches(y + 0.45), color, 1.5)
    add_text(slide, Inches(center_x + 0.06), Inches(center_y + 0.2), Inches(1.35), Inches(0.45), "六维评价\n框架", size=16, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
    add_panel(slide, Inches(7.58), Inches(1.92), Inches(4.72), Inches(4.4))
    add_text(slide, Inches(7.88), Inches(2.18), Inches(3.6), Inches(0.25), "判断链条", size=18, bold=True)
    stages = [
        ("起点", "问题创新性"),
        ("定位", "现状洞察度"),
        ("工具", "分析框架建构力"),
        ("骨架", "逻辑严密性"),
        ("落脚", "结论可接受性"),
        ("开放", "前瞻延展性"),
    ]
    y = 2.6
    for idx, (stage, dim) in enumerate(stages):
        stage_box = add_panel(slide, Inches(7.92), Inches(y), Inches(0.92), Inches(0.42), fill=RGBColor(245, 249, 253))
        stage_box.line.color.rgb = CYAN
        add_text(slide, Inches(8.02), Inches(y + 0.1), Inches(0.72), Inches(0.14), stage, size=11, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
        add_text(slide, Inches(9.0), Inches(y + 0.09), Inches(2.7), Inches(0.18), dim, size=13, color=TEXT)
        if idx < len(stages) - 1:
            connector(slide, Inches(8.38), Inches(y + 0.42), Inches(8.38), Inches(y + 0.52), CYAN, 1.1)
        y += 0.58
    add_panel(slide, Inches(7.92), Inches(5.92), Inches(4.0), Inches(0.56), fill=RGBColor(246, 250, 253))
    add_text(slide, Inches(8.1), Inches(6.08), Inches(3.65), Inches(0.22), "权重来自当前法学配置，六维共同构成完整评价路径。", size=11, color=MUTED, align=PP_ALIGN.CENTER)


def slide_iteration(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "这套机制不是一次成型，而是沿着四条方向持续收敛", "关键不是追版本，而是让评审尺度越来越可解释、可校准、可复核")
    tiles = [
        ("概念清单化", "把模糊概念改成检查清单，减少凭感觉打分。", BLUE),
        ("关键维度操作化", "研究地图、分析框架、结论可接受性、前瞻延展性逐步细化。", CYAN),
        ("封顶规则收窄", "持续收窄 ceiling rules 触发条件，避免轻易封顶。", GOLD),
        ("输出结构收短", "让输出更短、更结构化，更利于专家快速核验。", RGBColor(72, 120, 162)),
    ]
    positions = [(0.95, 2.0), (6.85, 2.0), (0.95, 4.15), (6.85, 4.15)]
    for (title, body, color), (x, y) in zip(tiles, positions):
        tile = add_panel(slide, Inches(x), Inches(y), Inches(5.1), Inches(1.65), fill=RGBColor(250, 252, 255))
        tile.line.color.rgb = color
        add_text(slide, Inches(x + 0.25), Inches(y + 0.18), Inches(2.6), Inches(0.25), title, size=18, bold=True, color=color)
        add_text(slide, Inches(x + 0.25), Inches(y + 0.58), Inches(4.5), Inches(0.6), body, size=14)
    center = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(5.95), Inches(3.18), Inches(1.42), Inches(0.82))
    center.fill.solid()
    center.fill.fore_color.rgb = RGBColor(240, 246, 250)
    center.line.color.rgb = BLUE
    add_text(slide, Inches(6.1), Inches(3.42), Inches(1.1), Inches(0.15), "持续收敛", size=14, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    for x1, y1, x2, y2 in [
        (5.95, 3.55, 5.35, 2.8), (7.37, 3.55, 7.95, 2.8), (5.95, 3.55, 5.35, 4.95), (7.37, 3.55, 7.95, 4.95),
    ]:
        connector(slide, Inches(x1), Inches(y1), Inches(x2), Inches(y2), CYAN, 1.3)


def slide_order(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "单维评分必须先抽证据，再分档给分", "固定顺序把印象打分改成证据驱动判断")
    steps = ["定位评价对象", "抽取关键证据", "判断上限规则", "确定分档给分", "输出摘要与理由"]
    coords = [(1.0, 3.8), (3.2, 2.5), (5.45, 3.85), (7.75, 2.55), (10.0, 3.9)]
    for idx, (label, (x, y)) in enumerate(zip(steps, coords), 1):
        card = add_panel(slide, Inches(x), Inches(y), Inches(1.8), Inches(1.0), fill=RGBColor(250, 252, 255))
        card.line.color.rgb = BLUE if idx % 2 else CYAN
        add_text(slide, Inches(x + 0.14), Inches(y + 0.12), Inches(0.4), Inches(0.2), f"{idx:02d}", size=12, bold=True, color=GOLD)
        add_text(slide, Inches(x + 0.14), Inches(y + 0.42), Inches(1.45), Inches(0.25), label, size=13, bold=True)
        if idx < len(steps):
            x2, y2 = coords[idx]
            connector(slide, Inches(x + 1.8), Inches(y + 0.5), Inches(x2), Inches(y2 + 0.5), CYAN, 1.4)
    add_panel(slide, Inches(0.95), Inches(5.55), Inches(11.4), Inches(0.72), fill=RGBColor(246, 250, 253))
    add_text(slide, Inches(1.15), Inches(5.78), Inches(11.0), Inches(0.2), "原则：禁止先凭整体印象给分，必须先定位对象、抽证据、判规则，再进入分档与分数。", size=15, bold=True, color=TEXT, align=PP_ALIGN.CENTER)


def slide_contract(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "统一输出契约让 AI 的判断可解释、可比较、可复核", "关键不是只给分，而是留下可供专家审核的判断结构")
    add_panel(slide, Inches(0.85), Inches(1.9), Inches(6.05), Inches(4.55))
    add_text(slide, Inches(1.1), Inches(2.16), Inches(2.8), Inches(0.25), "输出字段示意", size=18, bold=True, color=TEXT)
    fields = [
        "score / band", "summary / core_judgment", "score_rationale", "evidence_quotes", "strengths / weaknesses", "review_flags",
    ]
    x_vals = [1.12, 3.15]
    y = 2.78
    for idx, field in enumerate(fields):
        x = x_vals[idx % 2]
        yy = y + (idx // 2) * 0.88
        add_chip(slide, Inches(x), Inches(yy), Inches(1.65 if idx % 2 == 0 else 2.45), field, fill=RGBColor(245, 249, 253))
    add_panel(slide, Inches(7.3), Inches(2.1), Inches(4.8), Inches(3.9), fill=RGBColor(250, 252, 255))
    add_text(slide, Inches(7.55), Inches(2.36), Inches(3.5), Inches(0.25), "为什么要严格统一", size=18, bold=True)
    benefits = [
        ("可解释", "专家能看到 AI 为什么这么判。"),
        ("可比较", "不同模型可以横向比较。"),
        ("可复核", "事后可回看证据与规则。"),
    ]
    yy = 3.0
    for title, body in benefits:
        add_chip(slide, Inches(7.58), Inches(yy), Inches(1.0), title, fill=RGBColor(233, 242, 248))
        add_text(slide, Inches(8.76), Inches(yy + 0.08), Inches(2.9), Inches(0.25), body, size=14)
        yy += 0.86


def slide_reliability(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "可靠性判断决定哪里可以采纳，哪里必须请人介入", "重点不是谁分高，而是分歧有多大、结论可采纳到什么程度")
    add_panel(slide, Inches(0.95), Inches(2.05), Inches(5.35), Inches(3.9))
    add_text(slide, Inches(1.22), Inches(2.28), Inches(2.3), Inches(0.25), "模型得分示意", size=18, bold=True)
    bars = [("模型 A", 82, BLUE), ("模型 B", 79, CYAN), ("模型 C", 68, GOLD)]
    y = 2.95
    for label, score, color in bars:
        add_text(slide, Inches(1.22), Inches(y), Inches(0.78), Inches(0.2), label, size=13, bold=True)
        bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(2.08), Inches(y + 0.03), Inches(score / 18), Inches(0.22))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        add_text(slide, Inches(5.0), Inches(y - 0.02), Inches(0.6), Inches(0.2), str(score), size=12, color=MUTED)
        y += 0.62
    mean = add_panel(slide, Inches(1.2), Inches(4.95), Inches(1.35), Inches(0.7), fill=RGBColor(246, 250, 253))
    mean.line.color.rgb = BLUE
    add_text(slide, Inches(1.38), Inches(5.12), Inches(1.0), Inches(0.2), "均值 76.3", size=14, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    std = add_panel(slide, Inches(2.8), Inches(4.95), Inches(1.35), Inches(0.7), fill=RGBColor(246, 250, 253))
    std.line.color.rgb = GOLD
    add_text(slide, Inches(2.98), Inches(5.12), Inches(1.0), Inches(0.2), "std 7.4", size=14, bold=True, color=RGBColor(125, 94, 35), align=PP_ALIGN.CENTER)
    tag = add_panel(slide, Inches(4.4), Inches(4.95), Inches(1.45), Inches(0.7), fill=RGBColor(249, 243, 229))
    tag.line.color.rgb = GOLD
    add_text(slide, Inches(4.57), Inches(5.12), Inches(1.1), Inches(0.2), "进入复核", size=14, bold=True, color=RGBColor(125, 94, 35), align=PP_ALIGN.CENTER)

    add_panel(slide, Inches(7.1), Inches(2.0), Inches(5.2), Inches(3.95))
    add_text(slide, Inches(7.38), Inches(2.25), Inches(2.5), Inches(0.25), "判断意义", size=18, bold=True)
    add_bullets(slide, Inches(7.42), Inches(2.78), Inches(4.5), [
        "高置信度不等于无需专家，低置信度更明确提示必须复核。",
        "这一层把“模型输出”转成“可采纳程度判断”。",
        "当前阈值讨论本身也需要专家共同校准。",
    ], size=15, line_gap=0.48)


def slide_boundaries(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "当前边界仍然清楚存在，这正是专家意见最有价值的地方", "机制可以运作，但可信边界仍要靠共同体继续校准")
    # iceberg
    water = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.9), Inches(4.15), Inches(5.0), Inches(0.15))
    water.fill.solid()
    water.fill.fore_color.rgb = CYAN
    water.line.fill.background()
    poly_top = add_panel(slide, Inches(1.6), Inches(2.1), Inches(3.6), Inches(2.05), fill=RGBColor(248, 251, 254))
    poly_top.line.color.rgb = BLUE
    add_text(slide, Inches(2.2), Inches(2.85), Inches(2.4), Inches(0.35), "已结构化部分", size=20, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(2.05), Inches(3.28), Inches(2.7), Inches(0.3), "维度定义 / 决策顺序 / 输出契约", size=13, color=MUTED, align=PP_ALIGN.CENTER)
    poly_bottom = add_panel(slide, Inches(1.15), Inches(4.3), Inches(4.5), Inches(1.55), fill=RGBColor(237, 245, 249))
    poly_bottom.line.color.rgb = CYAN
    add_text(slide, Inches(2.0), Inches(4.8), Inches(2.8), Inches(0.4), "仍需专家判断部分", size=20, bold=True, color=CYAN, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(1.62), Inches(5.18), Inches(3.55), Inches(0.3), "维度边界 / 流派识别 / 证据充分性 / 封顶规则", size=13, color=MUTED, align=PP_ALIGN.CENTER)
    add_panel(slide, Inches(6.6), Inches(2.0), Inches(5.0), Inches(3.95))
    add_text(slide, Inches(6.9), Inches(2.25), Inches(2.6), Inches(0.25), "专家最值得发力的地方", size=18, bold=True)
    add_bullets(slide, Inches(6.95), Inches(2.8), Inches(4.3), [
        "哪些维度边界仍然模糊，需要更精确地切开。",
        "哪些规则触发过宽、过窄，导致 AI 失真或过度保守。",
        "哪些概念仍然过于抽象，尚未真正进入法学共同体语境。",
    ], size=15, line_gap=0.5)


def slide_roadmap(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "下一步分两段推进：1 个月校准，半年形成稳定机制", "目标不是更快自动化，而是把可用机制推进到可信机制")
    add_panel(slide, Inches(0.95), Inches(2.0), Inches(5.4), Inches(3.95), fill=RGBColor(250, 252, 255))
    add_text(slide, Inches(1.25), Inches(2.26), Inches(1.8), Inches(0.25), "未来 1 个月", size=20, bold=True, color=RGBColor(125, 94, 35))
    connector(slide, Inches(1.3), Inches(3.0), Inches(5.55), Inches(3.0), GOLD, 3)
    for x, title in [(1.55, "收集反馈"), (2.65, "校准规则"), (3.75, "补样本"), (4.85, "打磨复核")]:
        dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(2.83), Inches(0.28), Inches(0.28))
        dot.fill.solid()
        dot.fill.fore_color.rgb = GOLD
        dot.line.fill.background()
        add_text(slide, Inches(x - 0.18), Inches(3.18), Inches(0.65), Inches(0.22), title, size=11, color=TEXT, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(1.28), Inches(4.1), Inches(4.6), Inches(0.8), "核心目标：先把规则争议点暴露出来，让专家意见真正进入机制。", size=16)

    add_panel(slide, Inches(6.9), Inches(2.0), Inches(5.45), Inches(3.95), fill=RGBColor(250, 252, 255))
    add_text(slide, Inches(7.2), Inches(2.26), Inches(1.8), Inches(0.25), "未来半年", size=20, bold=True, color=BLUE)
    connector(slide, Inches(7.25), Inches(3.0), Inches(11.6), Inches(3.0), CYAN, 3)
    for x, title in [(7.55, "稳定规则"), (8.7, "标准样本库"), (10.0, "人机协同"), (11.1, "制度接入")]:
        dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(2.83), Inches(0.28), Inches(0.28))
        dot.fill.solid()
        dot.fill.fore_color.rgb = CYAN
        dot.line.fill.background()
        add_text(slide, Inches(x - 0.22), Inches(3.18), Inches(0.78), Inches(0.22), title, size=11, color=TEXT, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(7.22), Inches(4.1), Inches(4.7), Inches(0.8), "核心目标：形成更稳定的规则体系、样本体系和复核流程，再讨论跨学科扩展。", size=16)


def slide_closing(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_panel(slide, Inches(1.2), Inches(1.35), Inches(10.95), Inches(4.95), fill=RGBColor(248, 251, 254))
    add_text(slide, Inches(1.7), Inches(2.0), Inches(10.0), Inches(0.7), "希望专家帮助我们把机制从可用推进到可信", size=26, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
    add_chip(slide, Inches(2.2), Inches(3.25), Inches(2.6), "校准六维逻辑与边界规则")
    add_chip(slide, Inches(5.3), Inches(3.25), Inches(2.55), "参与样本试评与阶段性评估", fill=RGBColor(232, 247, 251), color=CYAN)
    add_chip(slide, Inches(8.35), Inches(3.25), Inches(2.3), "支持共建与试点推进", fill=RGBColor(249, 243, 229), color=RGBColor(125, 94, 35))
    add_text(slide, Inches(2.1), Inches(4.35), Inches(9.2), Inches(0.7), "如果这套方向成立，我们希望下一步不是“证明 AI 能替代评审”，而是与专家一起把评审标准做得更清楚、更可信。", size=18, color=MUTED, align=PP_ALIGN.CENTER)
    add_footer_note(slide, "期待专家对机制逻辑、规则边界、样本校准与试点路径提出意见。")


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_cover(prs)
    slide_problem(prs)
    slide_workflow(prs)
    slide_intent_gap(prs)
    slide_yaml(prs)
    slide_yaml_scope(prs)
    slide_yaml_rules(prs)
    slide_yaml_output_hooks(prs)
    slide_dimensions(prs)
    slide_iteration(prs)
    slide_order(prs)
    slide_contract(prs)
    slide_reliability(prs)
    slide_boundaries(prs)
    slide_roadmap(prs)
    slide_closing(prs)

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(str(PPTX_PATH))
    print(f"generated: {PPTX_PATH}")


if __name__ == "__main__":
    main()
